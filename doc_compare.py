"""Module to compare two documents using Difflib."""
import os
from os.path import abspath
from difflib import SequenceMatcher
import re
import shutil
import win32com.client as win32
import xlrd
import pptx
import csv
import sys
from datetime import datetime
base_path = os.path.dirname(abspath('__file__'))
results_path = base_path + '\\' + 'results'
csv_path = results_path + '\\' + 'results.csv'
if 'model' not in os.listdir(base_path):
    print('Add folder with model files. Folder should be named model.')
    sys.exit()
model_path = base_path + '\\' + 'model'
model_dict = {i.split('.')[0]:i.split('.')[1] for i in os.listdir(model_path)}


def get_jc(fname):
    """Get job code from file name."""
    jc = None
    parts = fname.split('_')
    if 'cn' in parts:
        if len(parts) == 4:
            jc = parts[0] + parts[1]
        elif len(parts) == 5:
            jc = parts[0] + parts[1] + parts[2]
    else:
        jc = fname.split('.')[0]
    return jc


def extract_text(fname, path=base_path):
    """Extract text from given document."""
    if fname.split('.')[-1] in ['doc', 'docx', 'rtf']:
        word = win32.Dispatch('Word.Application')
        doc = word.Documents.Open(path+'\\'+fname)
        txt = doc.Content.Text
        doc.Close(False)
    elif fname.split('.')[-1] in ['xls', 'xlsx']:
        workbook = xlrd.open_workbook(path+'\\'+fname)
        sheets_name = workbook.sheet_names()
        txt = '\n'
        for names in sheets_name:
            worksheet = workbook.sheet_by_name(names)
            num_rows = worksheet.nrows
            num_cells = worksheet.ncols
            for curr_row in range(num_rows):
                new_output = []
                for index_col in range(num_cells):
                    value = worksheet.cell_value(curr_row, index_col)
                    if value:
                        new_output.append(value)
                    if new_output:
                        txt += ' '.join([str(i) for i in new_output]) + '\n'
    elif fname.endswith('.pptx'):
        presentation = pptx.Presentation(path+'\\'+fname)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        txt = '\n\n'.join(text_runs)
    elif fname.endswith('.txt'):
        text_doc = open(path+'\\'+fname, 'r', encoding='utf8')
        txt = text_doc.read()
        text_doc.close()
    elif fname.endswith('.csv'):
        csv_doc = open(path+'\\'+fname, 'r', encoding='utf8')
        csv_reader = csv.reader(csv_doc, delimiter=',')
        txt = '\n'.join(['\t'.join(row) for row in csv_reader])
    return txt


def result_writer(result_arr):
    """Write results to csv file"""
    if 'results' not in os.listdir(base_path):
        os.mkdir('results')
    if 'results.csv' not in os.listdir(results_path):
        with open(csv_path, 'a', newline='') as result_csv:
            csv_writer = csv.writer(result_csv, delimiter=',')
            fields = ['file', 'solution_name', 'match_segments', 
                      'percent_match', 'percent_length_high']
            csv_writer.writerow(fields)
            result_csv.close()
    with open(csv_path, 'a', newline='') as result_csv:
        csv_writer = csv.writer(result_csv, delimiter=',')
        csv_writer.writerow(result_arr)
        result_csv.close()

    
for i in os.listdir(base_path):
    if i.startswith('test'):
        solution_name = i.split('_')[1]
        solution_path = base_path + '\\' + i
        for j in os.listdir(solution_path):
            f_name = j.split('.')[0]
            if f_name in model_dict.keys():
                job_code = get_jc(j)
                model_file = f_name + '.' + model_dict[f_name]
                model_text = extract_text(model_file, model_path)
                solution_text = extract_text(j, solution_path)
                similarity = SequenceMatcher(None, model_text, solution_text)
                matches = similarity.get_matching_blocks()
                percent = round(similarity.ratio() * 100, 2)
                high_matches = 0
                len_high_matches = 0
                for k in matches:
                    if k[2] > 70:
                        high_matches += 1
                        len_high_matches += k[2]
                test_doc_length = len(model_text) + len(solution_text)
                plhm = round(((2 * len_high_matches)/test_doc_length) * 100, 2)
                result_list = [job_code, solution_name, high_matches, 
                               percent, plhm]
                result_writer(result_list)
        print('Calculation complete for {}.'.format(i.split('_')[1]))
print('All calculations done.')
